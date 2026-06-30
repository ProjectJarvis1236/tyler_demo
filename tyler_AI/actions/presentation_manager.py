import json
import httpx
import random
import traceback
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

import configs

API_KEY = configs.OPENROUTER_KEY
URL = configs.URLS["nvidia/nemotron-3-super-120b-a12b:free"]

SYSTEM_PROMPT = """
Ты эксперт по созданию презентаций. На основе темы и количества слайдов (не обязательный параметр от пользователя, если отсутсвует, то составляй презентацию на 7 слайдов)
создай структуру презентации.

Верни СТРОГО JSON в формате:
{
    "slides": [
        {
            "type": "title",
            "title": "Заголовок презентации",
            "subtitle": "Подзаголовок или автор"
        },
        {
            "type": "content",
            "title": "Заголовок слайда",
            "bullets": ["Пункт 1", "Пункт 2", "Пункт 3"]
        },
        {
            "type": "two_column",
            "title": "Заголовок слайда",
            "left_column": ["Пункт 1", "Пункт 2"],
            "right_column": ["Пункт 1", "Пункт 2"]
        },
        {
            "type": "stats",
            "title": "Заголовок слайда",
            "stats": [
                {"number": "85%", "label": "Рост продаж"},
                {"number": "10K", "label": "Клиентов"},
                {"number": "24/7", "label": "Поддержка"}
            ]
        },
        {
            "type": "timeline",
            "title": "Заголовок слайда",
            "events": [
                {"year": "2020", "description": "Основание компании"},
                {"year": "2022", "description": "Выход на международный рынок"},
                {"year": "2024", "description": "Запуск нового продукта"}
            ]
        },
        {
            "type": "comparison",
            "title": "Заголовок слайда",
            "left_title": "Вариант А",
            "right_title": "Вариант Б",
            "left_items": ["Преимущество 1", "Преимущество 2"],
            "right_items": ["Недостаток 1", "Недостаток 2"]
        },
        {
            "type": "final",
            "title": "Спасибо за внимание",
            "subtitle": "Контакты или вывод"
        }
    ]
}

Типы слайдов:
title: титульный слайд (первый)
content: слайд с буллетами
two_column: слайд с двумя колонками
stats: слайд со статистикой (крупные цифры)
timeline: временная шкала (события по годам)
comparison: сравнение двух вариантов (заголовки + списки)
final: финальный слайд (последний)

Правила:
Первый слайд всегда type: "title"
Последний слайд всегда type: "final" и иметь фразу "Спасибо за вниание!"
Промежуточные — чередуй разные типы для разнообразия
Количество слайдов должно точно соответствовать запросу или содержать в себе от 2 - 20 слайдов
Заголовки и текст на русском языке или на языке, который укажет пользователь
Буллеты должны быть краткими и информативными
Используй stats для ключевых показателей
Используй timeline для истории, этапов, развития
Используй comparison для сравнений, плюсов/минусов. Каждый вариант должен содержать И плюсы, И минусы для честного сравнения. Добавляй маркеры в начале каждого пункта в полях left_items и right_items.
"""


def generate_random_gradient_colors():
    colors = [
        (31, 78, 121),
        (44, 62, 80),
        (52, 152, 219),
        (155, 89, 182),  # цвеиа
        (41, 128, 185),
        (22, 160, 133),
        (192, 57, 43),
        (230, 126, 34),
        (142, 68, 173),
        (26, 188, 156),
        (52, 73, 94),
        (241, 196, 15),
    ]

    count = random.randint(2, 4)
    colors_pre = random.sample(colors, count)
    return colors_pre


def apply_random_gradient(slide, colors_pre):  # для создания градиента в презентации
    num_colors = len(colors_pre)

    sld = slide._element

    cSld = sld.find(qn('p:cSld'))
    if cSld is None:
        cSld = etree.Element(qn('p:cSld'))
        sld.insert(0, cSld)

    for bg in cSld.findall(qn('p:bg')):
        cSld.remove(bg)

    bg = etree.Element(qn('p:bg'))
    cSld.insert(0, bg)

    bg_pr = etree.SubElement(bg, qn('p:bgPr'))

    grad_fill = etree.SubElement(bg_pr, qn('a:gradFill'))
    gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))

    for i, (r, g, b) in enumerate(colors_pre):
        pos = int(i * 100000 / (num_colors - 1))
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgb_clr = etree.SubElement(gs, qn('a:srgbClr'))
        srgb_clr.set('val', f"{r:02X}{g:02X}{b:02X}")

    lin = etree.SubElement(grad_fill, qn('a:lin'))
    lin.set('ang', '5400000')
    lin.set('scaled', '1')

    etree.SubElement(bg_pr, qn('a:effectLst'))


def extract_json(text):  # это нужно для того, чтобы отслеживать JSON структуру
    start = text.find('{')  # индекс 1-го вхождения
    if start == -1:
        return None

    balance = 0
    in_string = False
    escape_next = False

    for i in range(start, len(text)):
        char = text[i]

        if escape_next:
            escape_next = False
            continue

        if char == '\\':
            escape_next = True
            continue
        if char == '"':
            in_string = not in_string
            continue
        if not in_string:
            if char == '{':
                balance += 1
            elif char == '}':
                balance -= 1
                if balance == 0:
                    return text[start:i + 1]  # сплошной текст
    return None  # ничего нет


class PresentationManager:

    async def run(self, params: dict, chat_id: str):
        print("def_run")

        filename = params.get("filename", "presentation.pptx")
        if not filename.endswith('.pptx'):
            filename += '.pptx'

        folder = Path(configs.USER_FOLDER)  # думаю нужна ли тут проверка, что папка есть,
        # просто по идее прога в целом не запускается, если папки нет насколько я помню
        # (у меня была такая ошибка, как минимум)
        full_path = folder / filename

        topic = params.get("topic", "")
        slide_count = params.get("slide_count", 7)
        additional_info = params.get("additional_info", "")

        if not topic:
            return "не указана тема презинтации"

        slide_count = self.parse_slide_count(slide_count)
        print("количесвто слайдов", slide_count)

        structure = await self.generate_structure(topic, slide_count, additional_info)
        if not structure:
            return "Глеб должен мне буьылку сидра, нету генерации презы"  # { "status": "error", "message": "не сгенировать презинтацию"}

        result = self.create_presentation(full_path, structure)
        if not result:
            return "сохранение презентации не произошло"
        print(f"Презентация сохранена: {str(full_path)}")
        return result

    async def generate_structure(self, topic, slide_count, additional_info):
        prompt = {
            "model": "nvidia/nemotron-3-super-120b-a12b:free",
            "messages": [
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": f"""
    Создай структуру презентации:
    - Тема: {topic}
    - Количество слайдов: {slide_count} (включая титульный и финальный)
    - Дополнительная информация: {additional_info}

    Верни JSON структуру презентации согласно формату из системного промпта.
    """}
            ]
        }

        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }

        try:
            async with httpx.AsyncClient() as client:
                response = await client.post(URL, json=prompt, headers=headers, timeout=60.0)
                response.raise_for_status()  # ghjdthrf cnfnecf
                data = response.json()  # словарь после парсировки JSON

                content = data["choices"][0]["message"]["content"]

                json_str = extract_json(content)
                if json_str:
                    return json.loads(json_str)  # парсит строку JSON в Python-словарь
                else:
                    return json.loads(
                        content.strip())  # распарсить через весь ответ, если ллм умная и без болтливости все вернула хорошо

        except Exception as e:
            print(f"хуйня, переделывай {e}")
            traceback.print_exc()
            return None

    def create_presentation(self, full_path, structure):
        try:
            prs = Presentation()

            # параметры слайда
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            slides_data = structure.get("slides", [])
            if not slides_data:
                return None

            # Словарь для вызова нужного метода по типу слайда
            slide_methods = {
                "title": self.add_title_slide,
                "content": self.add_content_slide,
                "two_column": self.add_two_column_slide,
                "stats": self.add_stats_slide,
                "timeline": self.add_timeline_slide,
                "comparison": self.add_comparison_slide,
                "final": self.add_final_slide
            }

            for slide_data in slides_data:
                slide_type = slide_data.get("type")  # определение типа слайда
                method = slide_methods.get(slide_type)
                if method:
                    method(prs, slide_data)  # jnhbcjdrf
                else:
                    self.add_content_slide(prs, slide_data)  # просто обычный текстовый слайд

                # Применяем градиент ко всем слайдам
            gradient_colors = generate_random_gradient_colors()

            for slide in prs.slides:
                apply_random_gradient(slide, gradient_colors)
            prs.save(str(full_path))  # файл сохранен
            return str(full_path)

        except Exception as e:
            print(f" хуйня переделывай {e}", flush=True)
            traceback.print_exc()
            return None

    def add_stats_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # пустой слайд

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333),
                                             Inches(1))  # слева, сверху, ширина, высота
        title_frame = title_box.text_frame
        title_frame.word_wrap = True  # автоматический перенос слова,чтобы не было криво и не выходило за рамки
        title_frame.text = data.get("title", "")
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        stats = data.get("stats", [])
        if stats:
            num_stats = len(stats)
            box_width = 12.333 / num_stats

            for i, stat in enumerate(stats):
                left_pos = 0.5 + (i * box_width)
                number_box = slide.shapes.add_textbox(Inches(left_pos), Inches(2.5), Inches(box_width), Inches(1.5))
                number_frame = number_box.text_frame
                number_frame.word_wrap = True
                number_frame.text = stat.get("number", "")
                number_frame.paragraphs[0].font.size = Pt(24)
                number_frame.paragraphs[0].font.bold = True
                number_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # по центру

                label_box = slide.shapes.add_textbox(Inches(left_pos), Inches(4.5), Inches(box_width), Inches(1))
                label_frame = label_box.text_frame
                label_frame.word_wrap = True
                label_frame.text = stat.get("label", "")
                label_frame.paragraphs[0].font.size = Pt(18)
                label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def add_title_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = data.get("title", "Презентация")
        title_box.text_frame.paragraphs[0].font.size = Pt(54)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.333), Inches(1))
            sub_box.text_frame.word_wrap = True
            sub_box.text_frame.text = data.get("subtitle")
            sub_box.text_frame.paragraphs[0].font.size = Pt(28)
            sub_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
            sub_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_content_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = data.get("title", "")
        title_box.text_frame.paragraphs[0].font.size = Pt(40)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(4.5))
        content_box.text_frame.word_wrap = True

        bullets = data.get("bullets", [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = content_box.text_frame.paragraphs[0]
            else:
                p = content_box.text_frame.add_paragraph()

            p.text = f"• {bullet}"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(14)

    def add_final_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = data.get("title", "Спасибо за внимание!")
        title_box.text_frame.paragraphs[0].font.size = Pt(54)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.333), Inches(1))
            sub_box.text_frame.word_wrap = True
            sub_box.text_frame.text = data.get("subtitle")
            sub_box.text_frame.paragraphs[0].font.size = Pt(28)
            sub_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
            sub_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_two_column_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = data.get("title", "")
        title_box.text_frame.paragraphs[0].font.size = Pt(40)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(4.5))
        left_box.text_frame.word_wrap = True
        for i, item in enumerate(data.get("left_column", [])):
            if i == 0:
                p = left_box.text_frame.paragraphs[0]
            else:
                p = left_box.text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(10)

        right_box = slide.shapes.add_textbox(Inches(6.833), Inches(2), Inches(6), Inches(4.5))
        right_box.text_frame.word_wrap = True
        for i, item in enumerate(data.get("right_column", [])):
            if i == 0:
                p = right_box.text_frame.paragraphs[0]
            else:
                p = right_box.text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(10)

    def add_timeline_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = data.get("title", "")
        title_box.text_frame.paragraphs[0].font.size = Pt(40)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        events = data.get("events", [])
        if events:
            step = 12.333 / len(events)
            for i, event in enumerate(events):
                left = 0.5 + i * step
                year_box = slide.shapes.add_textbox(Inches(left), Inches(2.5), Inches(step), Inches(1))
                year_box.text_frame.word_wrap = True
                year_box.text_frame.text = str(event.get("year", ""))
                year_box.text_frame.paragraphs[0].font.size = Pt(28)
                year_box.text_frame.paragraphs[0].font.bold = True
                year_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                year_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                desc_box = slide.shapes.add_textbox(Inches(left), Inches(3.8), Inches(step), Inches(2))
                desc_box.text_frame.word_wrap = True
                desc_box.text_frame.text = event.get("description", "")
                desc_box.text_frame.paragraphs[0].font.size = Pt(18)
                desc_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(240, 240, 240)
                desc_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_comparison_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = data.get("title", "")
        title_box.text_frame.paragraphs[0].font.size = Pt(40)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # левая колонка, налево четкие пацаны не ходят
        ltitle = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(1))
        ltitle.text_frame.word_wrap = True
        ltitle.text_frame.text = data.get("left_title", "Вариант А")
        ltitle.text_frame.paragraphs[0].font.size = Pt(28)
        ltitle.text_frame.paragraphs[0].font.bold = True
        ltitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        ltitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        lbox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(6), Inches(3.5))
        lbox.text_frame.word_wrap = True
        for i, item in enumerate(data.get("left_items", [])):
            if i == 0:
                p = lbox.text_frame.paragraphs[0]
            else:
                p = lbox.text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # правая колонка
        rtitle = slide.shapes.add_textbox(Inches(6.833), Inches(2), Inches(6), Inches(1))
        rtitle.text_frame.word_wrap = True
        rtitle.text_frame.text = data.get("right_title", "Вариант Б")
        rtitle.text_frame.paragraphs[0].font.size = Pt(28)
        rtitle.text_frame.paragraphs[0].font.bold = True
        rtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        rtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        rbox = slide.shapes.add_textbox(Inches(6.833), Inches(3), Inches(6), Inches(3.5))
        rbox.text_frame.word_wrap = True
        for i, item in enumerate(data.get("right_items", [])):
            if i == 0:
                p = rbox.text_frame.paragraphs[0]
            else:
                p = rbox.text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)

    def parse_slide_count(self, slide_count):
        if isinstance(slide_count, int):
            return max(2, min(slide_count, 20))

        if isinstance(slide_count, str):
            slide_count = slide_count.strip().lower()
            # да костыли, не могу понять как по-другому реализовать
            # хочу различное количество слайдов, но не хочу указывать строго количество
            if slide_count in ["короткая", "короткую", "небольшая", "маленькая"]:
                return random.randint(2, 5)
            elif slide_count in ["средняя", "обычная", "стандартная"]:
                return random.randint(5, 10)
            elif slide_count in ["подробная", "подробную", "развёрнутая", "большая", "большую"]:
                return random.randint(10, 20)

            numbers = re.findall(r"\d+", slide_count)
            if numbers:
                count = int(numbers[0])
                return max(2, min(count, 20))

        return 7


# это шаблоны
"""    def __init__(self, template_path=None):
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()

    def add_slide(self, slide_data):
        slide_layout = self.prs.slide_layouts[1]  # Title + Content (как в шаблоне)
        slide = self.prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data.get("title", "Без названия")  # Заголовок

        top = Inches(1.5)
        items = slide_data.get("items")  # Буллеты (текст слайда)
        placeholder = slide.placeholders[1]  # для всех bullets на слайде
        tf = placeholder.text_frame
        tf.clear()
        tf.word_wrap = True

        for item in items:
            bullet = item.get("bullet")
            para = item.get("paragraph")

            if bullet:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(24)
                p.level = 0

            if para:
                p_para = tf.add_paragraph()
                p_para.text = para
                p_para.font.size = Pt(18)
                p_para.level = 1

        charts = slide_data.get("charts")  # TODO: диаграммы
        if charts:
            for chart in charts: pass  # chart = {"type": "bar", "data": {...}, "position": (left, top, width, height)}

        images = slide_data.get("images")
        if images:
            for img in images:
                slide.shapes.add_picture(
                    img["path"],
                    img.get("left", Inches(1)),
                    img.get("top", Inches(1)),
                    width=img.get("width", Inches(4)),
                    height=img.get("height", Inches(3))
                )

    def _add_title_slide(self, my_title):
        title_slide = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide)
        slide.shapes.title.text = my_title

"""
