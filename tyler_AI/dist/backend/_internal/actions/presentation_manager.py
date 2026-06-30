from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import asyncio



class PresentationManager:
    def __init__(self, template_path=None):
        if template_path: self.prs = Presentation(template_path)
        else: self.prs = Presentation()


    def add_slide(self, slide_data):
        slide_layout = self.prs.slide_layouts[1]    # Title + Content (как в шаблоне)
        slide = self.prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data.get("title", "Без названия") # Заголовок

        top = Inches(1.5)
        items = slide_data.get("items")     # Буллеты (текст слайда)
        placeholder = slide.placeholders[1]  # для всех bullets на слайде
        tf = placeholder.text_frame
        tf.clear()
        tf.word_wrap = True


        for item in items:
            bullet = item.get("bullet")
            para = item.get("paragraph")

            if bullet:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(24)
                p.level = 0

            if para:
                p_para = tf.add_paragraph()
                p_para.text = para
                p_para.font.size = Pt(18)
                p_para.level = 1


        charts = slide_data.get("charts")       # TODO: диаграммы
        if charts:
            for chart in charts: pass       # chart = {"type": "bar", "data": {...}, "position": (left, top, width, height)}


        images = slide_data.get("images")
        if images:
            for img in images:
                slide.shapes.add_picture(
                    img["path"],
                    img.get("left", Inches(1)),
                    img.get("top", Inches(1)),
                    width=img.get("width", Inches(4)),
                    height=img.get("height", Inches(3))
                )


    def _add_title_slide(self, my_title):
        title_slide = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide)
        slide.shapes.title.text = my_title


    async def run(self, params: dict, chat_id: str):
        title = params.get("title")
        if title:
            self._add_title_slide(title)

        slides = params.get("slides", [])
        filename = params.get("filename", "presentation.pptx")
        for slide_data in slides: self.add_slide(slide_data)

        loop = asyncio.get_running_loop()
        await loop.run_in_executor(None, self.prs.save, filename)
        print(f"Презентация сохранена: {filename}")


#Пока что так, без картинок, без диаграмм
"""

 {{
    "filename": "<Название файла с презентацией>"
    "title": "<Заголвок презентации (титульный слайд)>",
    "slides": [
        {{"title": "<Заголовок слайда>", "items": [ {{"bullet": "<Буллет1>", "paragraph": "<Текстовый блок>"}}, {{"bullets: "<Буллет1>", "paragraph": "<Текстовый блок>"}} ...],
        {{"title": "<Заголовок слайда>", "items": [ {{"bullet": "<Буллет1>", "paragraph": "<Текстовый блок>"}}, {{"bullets: "<Буллет1>", "paragraph": "<Текстовый блок>"}} ...],
        ...
    ],
}}

"""
        